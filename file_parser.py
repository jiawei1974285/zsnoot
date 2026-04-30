#!/usr/bin/env python3
"""
Multi-format file parser with explicit parse errors.
"""
import os
import subprocess
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path


class FileParseError(RuntimeError):
    """Raised when a file cannot be parsed into usable text."""


def parse_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".md":
        return parse_markdown(file_path)
    if ext == ".txt":
        return parse_text(file_path)
    if ext in [".docx", ".doc"]:
        return parse_docx(file_path)
    if ext in [".pptx", ".ppt"]:
        return parse_pptx(file_path)
    if ext == ".pdf":
        return parse_pdf(file_path)
    if ext in [".xlsx", ".xls"]:
        return parse_excel(file_path)
    if ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
        return parse_image(file_path)
    return parse_text(file_path)


def _ensure_nonempty(text: str, file_type: str) -> str:
    cleaned = (text or "").strip()
    if not cleaned:
        raise FileParseError(f"{file_type} 解析后内容为空")
    return cleaned


def parse_markdown(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return _ensure_nonempty(f.read(), "Markdown")


def parse_text(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return _ensure_nonempty(f.read(), "文本")
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="gbk") as f:
            return _ensure_nonempty(f.read(), "文本")


def parse_docx(file_path: str) -> str:
    try:
        from docx import Document

        doc = Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
        return _ensure_nonempty(text, "Word")
    except ImportError:
        return parse_docx_fallback(file_path)
    except Exception as exc:
        raise FileParseError(f"Word 解析失败: {exc}") from exc


def parse_docx_fallback(file_path: str) -> str:
    if Path(file_path).suffix.lower() == ".doc":
        raise FileParseError("DOC 需要 Word 解析组件，当前环境暂不支持")

    try:
        with zipfile.ZipFile(file_path) as archive:
            xml_bytes = archive.read("word/document.xml")
    except Exception as exc:
        raise FileParseError(f"Word 兜底解析失败: {exc}") from exc

    try:
        root = ET.fromstring(xml_bytes)
        namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        texts = [node.text for node in root.findall(".//w:t", namespaces) if node.text]
        return _ensure_nonempty("".join(texts), "Word")
    except Exception as exc:
        raise FileParseError(f"Word XML 解析失败: {exc}") from exc


def parse_pptx(file_path: str) -> str:
    if Path(file_path).suffix.lower() == ".ppt":
        raise FileParseError("PPT 需要 PowerPoint 解析组件，当前环境暂不支持，请转换为 PPTX")

    try:
        from pptx import Presentation

        presentation = Presentation(file_path)
        chunks = []
        for index, slide in enumerate(presentation.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                chunks.append(f"## Slide {index}\n" + "\n".join(texts))
        return _ensure_nonempty("\n\n".join(chunks), "PowerPoint")
    except ImportError:
        return parse_pptx_fallback(file_path)
    except Exception as exc:
        raise FileParseError(f"PowerPoint 解析失败: {exc}") from exc


def parse_pptx_fallback(file_path: str) -> str:
    """Extract slide text from PPTX XML without external dependencies."""
    try:
        with zipfile.ZipFile(file_path) as archive:
            slide_names = sorted(
                (
                    name
                    for name in archive.namelist()
                    if re_match_pptx_slide(name)
                ),
                key=pptx_slide_sort_key,
            )
            chunks = []
            for index, slide_name in enumerate(slide_names, 1):
                root = ET.fromstring(archive.read(slide_name))
                texts = [
                    node.text.strip()
                    for node in root.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t")
                    if node.text and node.text.strip()
                ]
                if texts:
                    chunks.append(f"## Slide {index}\n" + "\n".join(texts))
        return _ensure_nonempty("\n\n".join(chunks), "PowerPoint")
    except Exception as exc:
        raise FileParseError(f"PowerPoint 兜底解析失败: {exc}") from exc


def re_match_pptx_slide(name: str) -> bool:
    import re

    return bool(re.fullmatch(r"ppt/slides/slide\d+\.xml", name))


def pptx_slide_sort_key(name: str) -> int:
    import re

    match = re.search(r"slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else 0


def parse_pdf(file_path: str) -> str:
    try:
        import fitz

        doc = fitz.open(file_path)
        try:
            text = "\n".join(page.get_text() for page in doc)
        finally:
            doc.close()
        return _ensure_nonempty(text, "PDF")
    except ImportError:
        pass
    except Exception as exc:
        raise FileParseError(f"PDF 解析失败: {exc}") from exc

    try:
        result = subprocess.run(["pdftotext", file_path, "-"], capture_output=True, text=True, check=False)
    except FileNotFoundError as exc:
        raise FileParseError("PDF 解析依赖缺失，请安装 PyMuPDF 或 pdftotext") from exc

    if result.returncode != 0:
        raise FileParseError(result.stderr.strip() or "PDF 解析命令执行失败")
    return _ensure_nonempty(result.stdout, "PDF")


def parse_excel(file_path: str) -> str:
    try:
        import pandas as pd

        excel = pd.ExcelFile(file_path)
        chunks = []
        for sheet_name in excel.sheet_names:
            sheet = excel.parse(sheet_name)
            chunks.append(f"## {sheet_name}\n{sheet.to_string(index=False)}")
        return _ensure_nonempty("\n\n".join(chunks), "Excel")
    except ImportError:
        return parse_excel_fallback(file_path)
    except Exception as exc:
        raise FileParseError(f"Excel 解析失败: {exc}") from exc


def parse_excel_fallback(file_path: str) -> str:
    if Path(file_path).suffix.lower() == ".xls":
        raise FileParseError("XLS 需要 Excel 解析组件，当前环境暂不支持")

    try:
        with zipfile.ZipFile(file_path) as archive:
            shared_strings = _read_shared_strings(archive)
            workbook = ET.fromstring(archive.read("xl/workbook.xml"))
            rels = ET.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
            ns = {
                "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
                "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
            }
            rel_map = {
                rel.attrib["Id"]: rel.attrib["Target"]
                for rel in rels.findall(".//rel:Relationship", ns)
            }

            chunks = []
            for sheet in workbook.findall(".//main:sheets/main:sheet", ns):
                name = sheet.attrib.get("name", "Sheet")
                rel_id = sheet.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                target = rel_map.get(rel_id, "")
                if not target:
                    continue
                xml_path = f"xl/{target}" if not target.startswith("xl/") else target
                sheet_xml = ET.fromstring(archive.read(xml_path))
                rows = []
                for row in sheet_xml.findall(".//main:sheetData/main:row", ns):
                    cells = []
                    for cell in row.findall("main:c", ns):
                        value = _read_excel_cell(cell, ns, shared_strings)
                        if value is not None:
                            cells.append(str(value))
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    chunks.append(f"## {name}\n" + "\n".join(rows))
        return _ensure_nonempty("\n\n".join(chunks), "Excel")
    except Exception as exc:
        raise FileParseError(f"Excel 兜底解析失败: {exc}") from exc


def _read_shared_strings(archive: zipfile.ZipFile):
    try:
        root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
    except KeyError:
        return []
    ns = {"main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    values = []
    for item in root.findall(".//main:si", ns):
        text = "".join(node.text or "" for node in item.findall(".//main:t", ns))
        values.append(text)
    return values


def _read_excel_cell(cell, ns, shared_strings):
    cell_type = cell.attrib.get("t")
    if cell_type == "inlineStr":
        text_node = cell.find("main:is/main:t", ns)
        return text_node.text if text_node is not None else None
    value_node = cell.find("main:v", ns)
    if value_node is None:
        return None
    raw = value_node.text or ""
    if cell_type == "s":
        index = int(raw)
        return shared_strings[index] if 0 <= index < len(shared_strings) else raw
    return raw


def parse_image(file_path: str) -> str:
    try:
        import pytesseract
        from PIL import Image

        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang="chi_sim+eng")
        return _ensure_nonempty(text, "图片 OCR")
    except ImportError as exc:
        raise FileParseError("图片 OCR 依赖缺失，请安装 Pillow、pytesseract 和 Tesseract") from exc
    except Exception as exc:
        raise FileParseError(f"图片 OCR 失败: {exc}") from exc


def get_file_info(file_path: str) -> dict:
    stat = os.stat(file_path)
    return {
        "name": os.path.basename(file_path),
        "path": file_path,
        "size": stat.st_size,
        "ext": Path(file_path).suffix.lower(),
        "created": stat.st_ctime,
        "modified": stat.st_mtime,
    }
